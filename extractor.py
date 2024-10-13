from abc import ABC, abstractmethod
import os
import csv
from PIL import Image
from io import BytesIO
from pdfminer.high_level import extract_text
import pdfplumber
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

# Abstract Class for File Loading
class FileLoader(ABC):
    @abstractmethod
    def validate_file(self) -> bool:
        pass

    @abstractmethod
    def load_file(self):
        pass

# Concrete PDFLoader Class
class PDFLoader(FileLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def validate_file(self) -> bool:
        return self.file_path.endswith('.pdf') and os.path.isfile(self.file_path)

    def load_file(self):
        if not self.validate_file():
            raise ValueError("Invalid PDF file")
        
        # Try to extract text with pdfminer
        try:
            text = extract_text(self.file_path)
            if text.strip():
                return text
            else:
                print("No text found in PDF using pdfminer, falling back to OCR.")
        except Exception as e:
            print(f"Error extracting text with pdfminer: {e}, falling back to OCR.")
        
        # Fallback to OCR if no text is found
        try:
            return self.extract_text_with_ocr()
        except Exception as e:
            print(f"Error during OCR extraction: {e}")
            raise ValueError("Failed to extract text from PDF.")

    def extract_text_with_ocr(self):
        """Convert PDF pages to images and perform OCR."""
        images = convert_from_path(self.file_path)
        text = ""
        for image in images:
            try:
                ocr_text = pytesseract.image_to_string(image)
                text += ocr_text
            except Exception as e:
                print(f"Error during OCR on one of the pages: {e}")
        return text

    def extract_links(self):
        """Extract hyperlinks from the PDF."""
        links = []
        with pdfplumber.open(self.file_path) as pdf:
            for page in pdf.pages:
                if page.annots:  # Check if annotations are available
                    for annotation in page.annots:
                        if annotation.get("uri"):
                            links.append(annotation["uri"])
        return links

    def extract_images(self):
        """Extract images from the PDF."""
        images = []
        with pdfplumber.open(self.file_path) as pdf:
            for page in pdf.pages:
                if page.images:
                    images.extend(page.images)  # Append images metadata
        return images

    def extract_tables(self):
        """Extract tables from the PDF."""
        tables = []
        with pdfplumber.open(self.file_path) as pdf:
            for page in pdf.pages:
                tables.extend(page.extract_tables())  # Append all tables
        return tables

# Concrete DOCXLoader Class
class DOCXLoader(FileLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.content = None

    def validate_file(self) -> bool:
        return self.file_path.endswith('.docx') and os.path.isfile(self.file_path)

    def load_file(self):
        if self.validate_file():
            self.content = Document(self.file_path)
            return self.content
        raise ValueError("Invalid DOCX file")

# Concrete PPTLoader Class
class PPTLoader(FileLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.content = None

    def validate_file(self) -> bool:
        return self.file_path.endswith('.pptx') and os.path.isfile(self.file_path)

    def load_file(self):
        if self.validate_file():
            self.content = Presentation(self.file_path)
            return self.content
        raise ValueError("Invalid PPT file")

# DataExtractor Class
class DataExtractor:
    def __init__(self, file_loader: FileLoader):
        self.file_loader = file_loader
        self.content = self.file_loader.load_file()

    def extract_text(self):
        if isinstance(self.file_loader, PDFLoader):
            return self.content
        elif isinstance(self.file_loader, DOCXLoader):
            return '\n'.join(paragraph.text for paragraph in self.content.paragraphs)
        elif isinstance(self.file_loader, PPTLoader):
            return '\n'.join(
                shape.text for slide in self.content.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        return ""

    def extract_links(self):
        if isinstance(self.file_loader, PDFLoader):
            return self.file_loader.extract_links()
        links = []
        if isinstance(self.file_loader, DOCXLoader):
            for rel in self.content.part.rels.values():
                if rel.reltype == RT.HYPERLINK:
                    links.append(rel._target)
        elif isinstance(self.file_loader, PPTLoader):
            for slide in self.content.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    links.append((run.text, run.hyperlink.address))
        return links

    def extract_images(self):
        if isinstance(self.file_loader, PDFLoader):
            return self.file_loader.extract_images()
        images = []
        if isinstance(self.file_loader, DOCXLoader):
            for rel in self.content.part.rels.values():
                if "image" in rel.reltype:
                    images.append(rel.target_ref)
        elif isinstance(self.file_loader, PPTLoader):
            for slide in self.content.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        images.append(shape.image.blob)
        return images

    def extract_tables(self):
        if isinstance(self.file_loader, PDFLoader):
            return self.file_loader.extract_tables()
        tables = []
        if isinstance(self.file_loader, DOCXLoader):
            for table in self.content.tables:
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                tables.append(table_data)
        elif isinstance(self.file_loader, PPTLoader):
            for slide in self.content.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                        tables.append(table_data)
        return tables

# Storage class to handle saving extracted data
class Storage:
    def __init__(self, extractor: DataExtractor, base_path: str):
        self.extractor = extractor
        self.base_path = base_path
        self.ensure_directories_exist()

    def ensure_directories_exist(self):
        os.makedirs(os.path.join(self.base_path, 'images'), exist_ok=True)
        os.makedirs(os.path.join(self.base_path, 'tables'), exist_ok=True)
        os.makedirs(os.path.join(self.base_path, 'text'), exist_ok=True)
        os.makedirs(os.path.join(self.base_path, 'links'), exist_ok=True)

    def save_text(self):
        text = self.extractor.extract_text()
        file_type = self._get_file_type()
        text_file = os.path.join(self.base_path, 'text', f'{file_type}_text.txt')
        try:
            with open(text_file, 'w', encoding='utf-8') as f:
                f.write(text)
            print(f"Successfully saved text to {text_file}")
        except Exception as e:
            print(f"Error saving text: {e}")

    def save_links(self):
        links = self.extractor.extract_links()
        file_type = self._get_file_type()
        links_file = os.path.join(self.base_path, 'links', f'{file_type}_links.txt')
        try:
            with open(links_file, 'w', encoding='utf-8') as f:
                for link in links:
                    f.write(f"{link}\n")
            print(f"Successfully saved links to {links_file}")
        except Exception as e:
            print(f"Error saving links: {e}")

    def save_images(self):
        images = self.extractor.extract_images()
        image_folder = os.path.join(self.base_path, 'images')
        
        # Determine the file type prefix
        file_type_prefix = self._get_file_type()  # This should return 'pdf', 'docx', or 'ppt'

        for idx, image_data in enumerate(images):
            try:
                if isinstance(image_data, dict):  # For pdfplumber images (metadata)
                    image_data = BytesIO(image_data['stream'].get_data())
                elif isinstance(image_data, bytes):  # For DOCX/PPT images (binary)
                    image_data = BytesIO(image_data)
                
                # Open the image using PIL
                image = Image.open(image_data)

                # Construct image path with appropriate prefix and extension
                image_path = os.path.join(image_folder, f'{file_type_prefix}_image_{idx + 1}.{image.format.lower()}')

                # Save the image
                image.save(image_path)
                print(f"Successfully saved image {idx + 1} to {image_path}")
            except Exception as e:
                print(f"Error saving image {idx + 1}: {e}")

    def save_tables(self):
        tables = self.extractor.extract_tables()
        file_type = self._get_file_type()
        table_folder = os.path.join(self.base_path, 'tables')
        for idx, table in enumerate(tables):
            csv_path = os.path.join(table_folder, f'table_{file_type}_{idx + 1}.csv')
            try:
                with open(csv_path, mode='w', newline='', encoding='utf-8') as file:
                    writer = csv.writer(file)
                    writer.writerows(table)
                print(f"Successfully saved table {idx + 1} to {csv_path}")
            except Exception as e:
                print(f"Error saving table {idx + 1}: {e}")

    def _get_file_type(self):
        """Helper method to get the file type (pdf, docx, ppt) based on the loader class."""
        if isinstance(self.extractor.file_loader, PDFLoader):
            return 'pdf'
        elif isinstance(self.extractor.file_loader, DOCXLoader):
            return 'docx'
        elif isinstance(self.extractor.file_loader, PPTLoader):
            return 'ppt'
        return 'unknown'

# Example Usage
if __name__ == "__main__":
    # Specify file paths and base output folder
    pdf_file = 'sample.pdf'
    docx_file = 'sample.docx'
    ppt_file = 'sample.pptx'
    base_output_folder = 'output_data'

    # Extracting data from PDF and saving it
    pdf_loader = PDFLoader(pdf_file)
    pdf_extractor = DataExtractor(pdf_loader)
    pdf_storage = Storage(pdf_extractor, base_output_folder)
    pdf_storage.save_text()
    pdf_storage.save_links()
    pdf_storage.save_images()
    pdf_storage.save_tables()

    # Extracting data from DOCX and saving it
    docx_loader = DOCXLoader(docx_file)
    docx_extractor = DataExtractor(docx_loader)
    docx_storage = Storage(docx_extractor, base_output_folder)
    docx_storage.save_text()
    docx_storage.save_links()
    docx_storage.save_images()
    docx_storage.save_tables()

    # Extracting data from PPTX and saving it
    ppt_loader = PPTLoader(ppt_file)
    ppt_extractor = DataExtractor(ppt_loader)
    ppt_storage = Storage(ppt_extractor, base_output_folder)
    ppt_storage.save_text()
    ppt_storage.save_links()
    ppt_storage.save_images()
    ppt_storage.save_tables()
